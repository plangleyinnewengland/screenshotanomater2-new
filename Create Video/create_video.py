"""
Create a narrated video from a PowerPoint presentation.
- Extracts embedded images from each slide
- Uses external screenshot files referenced in speaker notes when available
- Generates TTS narration from speaker notes using edge-tts
- Assembles into an MP4 video with moviepy
"""
import os
import re
import asyncio
import threading
import tempfile
import shutil
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# Default configuration values
DEFAULTS = {
    "pptx_path": r"C:\GitHub\PDC\Monitor Your Services\Monitor Services.pptx",
    "output_dir": r"C:\GitHub\PDC\Monitor Your Services\video_build",
    "output_video": r"C:\GitHub\PDC\Monitor Your Services\Monitor Services.mp4",
    "voice": "en-US-GuyNeural",
}


def extract_slides(pptx_path):
    """Extract image and narration data from each slide."""
    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        # Get speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # Check for external file path in notes (first line or lines ending in .png/.jpg)
        external_path = None
        narration_lines = []
        for line in notes.split("\n"):
            line = line.strip()
            if not line:
                continue
            # Check if this line is a file path (drive letter with / or \, UNC, or relative paths)
            if re.match(r'^([A-Za-z]:[/\\]|\\\\|//).*\.(png|jpg|jpeg|gif|bmp)$', line, re.IGNORECASE):
                external_path = line
            else:
                narration_lines.append(line)

        narration = " ".join(narration_lines).strip()

        # Extract the main embedded image from the slide
        embedded_img = None
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                img_bytes = shape.image.blob
                embedded_img = img_bytes
                break  # Take the first/main image

        slides.append({
            'slide_num': i,
            'external_path': external_path,
            'embedded_img': embedded_img,
            'narration': narration,
        })

    return slides


def save_slide_images(slides, output_dir):
    """Save the best available image for each slide."""
    image_paths = []
    for s in slides:
        img_path = os.path.join(output_dir, f"slide_{s['slide_num']:03d}.png")

        # Prefer external screenshot if it exists
        if s['external_path'] and os.path.exists(s['external_path']):
            img = Image.open(s['external_path'])
            img.save(img_path, "PNG")
            print(f"  Slide {s['slide_num']}: Using external screenshot")
        elif s['embedded_img']:
            img = Image.open(BytesIO(s['embedded_img']))
            img.save(img_path, "PNG")
            print(f"  Slide {s['slide_num']}: Using embedded image")
        else:
            # Create a blank slide placeholder
            img = Image.new('RGB', (1920, 1080), (30, 30, 30))
            img.save(img_path, "PNG")
            print(f"  Slide {s['slide_num']}: Using blank placeholder")

        image_paths.append(img_path)

    return image_paths


async def generate_tts(slides, output_dir, voice):
    """Generate TTS audio for each slide with narration."""
    import edge_tts

    audio_paths = []
    durations = []

    for s in slides:
        audio_path = os.path.join(output_dir, f"slide_{s['slide_num']:03d}.mp3")

        if s['narration']:
            communicate = edge_tts.Communicate(s['narration'], voice)
            await communicate.save(audio_path)
            print(f"  Slide {s['slide_num']}: Generated audio ({len(s['narration'])} chars)")
            audio_paths.append(audio_path)
            durations.append(None)  # Will determine from audio file
        else:
            # No narration - use a short display duration
            audio_paths.append(None)
            durations.append(3.0)  # 3 second pause for slides without narration
            print(f"  Slide {s['slide_num']}: No narration (3s pause)")

    return audio_paths, durations


def assemble_video(image_paths, audio_paths, durations, output_dir, output_video):
    """Combine images and audio into final video."""
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips

    # Target resolution
    TARGET_W, TARGET_H = 1920, 1080

    clips = []
    for i, (img_path, audio_path, fallback_dur) in enumerate(zip(image_paths, audio_paths, durations)):
        # Load and resize image to target resolution
        img = Image.open(img_path)

        # Calculate scaling to fit within target while maintaining aspect ratio
        img_w, img_h = img.size
        scale = min(TARGET_W / img_w, TARGET_H / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)

        # Create black background and paste centered image
        canvas = Image.new('RGB', (TARGET_W, TARGET_H), (0, 0, 0))
        resized = img.resize((new_w, new_h), Image.LANCZOS)
        x_offset = (TARGET_W - new_w) // 2
        y_offset = (TARGET_H - new_h) // 2
        canvas.paste(resized, (x_offset, y_offset))

        # Save the composited frame
        frame_path = os.path.join(output_dir, f"frame_{i:03d}.png")
        canvas.save(frame_path, "PNG")

        if audio_path and os.path.exists(audio_path):
            audio_clip = AudioFileClip(audio_path)
            duration = audio_clip.duration + 1.0  # Add 1s buffer after narration
            img_clip = ImageClip(frame_path).with_duration(duration)
            img_clip = img_clip.with_audio(audio_clip)
        else:
            duration = fallback_dur or 3.0
            img_clip = ImageClip(frame_path).with_duration(duration)

        clips.append(img_clip)
        slide_num = i + 1
        print(f"  Slide {slide_num}: {duration:.1f}s")

    print(f"\nConcatenating {len(clips)} clips...")
    final = concatenate_videoclips(clips, method="compose")

    print(f"Writing video to: {output_video}")
    final.write_videofile(
        output_video,
        fps=24,
        codec="libx264",
        audio_codec="aac",
        preset="medium",
        threads=4,
    )
    print(f"\nVideo created: {output_video}")
    print(f"Duration: {final.duration:.1f}s ({final.duration/60:.1f} min)")


async def run_pipeline(pptx_path, output_dir, output_video, voice, log_fn):
    """Run the full video creation pipeline, reporting progress via log_fn."""
    log_fn("=" * 50)
    log_fn("VIDEO CREATOR")
    log_fn("=" * 50)

    log_fn("\n[1/4] Extracting slides from PowerPoint...")
    slides = extract_slides(pptx_path)
    log_fn(f"  Found {len(slides)} slides")

    log_fn("\n[2/4] Saving slide images...")
    image_paths = save_slide_images(slides, output_dir)

    log_fn("\n[3/4] Generating TTS narration...")
    audio_paths, durations = await generate_tts(slides, output_dir, voice)

    log_fn("\n[4/4] Assembling video...")
    assemble_video(image_paths, audio_paths, durations, output_dir, output_video)

    log_fn("\nDone!")


class VideoCreatorApp:
    """Tkinter GUI for the video creator."""

    def __init__(self, root):
        self.root = root
        root.title("Video Creator")
        root.resizable(False, False)

        # --- Configuration frame ---
        config_frame = ttk.LabelFrame(root, text="Configuration", padding=10)
        config_frame.grid(row=0, column=0, padx=10, pady=(10, 5), sticky="ew")

        # PowerPoint file
        ttk.Label(config_frame, text="PowerPoint file:").grid(row=0, column=0, sticky="w", pady=2)
        self.pptx_var = tk.StringVar(value=DEFAULTS["pptx_path"])
        ttk.Entry(config_frame, textvariable=self.pptx_var, width=70).grid(row=0, column=1, padx=5, pady=2)
        ttk.Button(config_frame, text="Browse...", command=self._browse_pptx).grid(row=0, column=2, pady=2)

        # Output directory
        ttk.Label(config_frame, text="Build directory:").grid(row=1, column=0, sticky="w", pady=2)
        self.outdir_var = tk.StringVar(value=DEFAULTS["output_dir"])
        ttk.Entry(config_frame, textvariable=self.outdir_var, width=70).grid(row=1, column=1, padx=5, pady=2)
        ttk.Button(config_frame, text="Browse...", command=self._browse_outdir).grid(row=1, column=2, pady=2)

        # Output video
        ttk.Label(config_frame, text="Output video:").grid(row=2, column=0, sticky="w", pady=2)
        self.video_var = tk.StringVar(value=DEFAULTS["output_video"])
        ttk.Entry(config_frame, textvariable=self.video_var, width=70).grid(row=2, column=1, padx=5, pady=2)
        ttk.Button(config_frame, text="Save As...", command=self._browse_video).grid(row=2, column=2, pady=2)

        # TTS Voice
        ttk.Label(config_frame, text="TTS Voice:").grid(row=3, column=0, sticky="w", pady=2)
        self.voice_var = tk.StringVar(value=DEFAULTS["voice"])
        ttk.Entry(config_frame, textvariable=self.voice_var, width=70).grid(row=3, column=1, padx=5, pady=2)

        # --- Buttons ---
        btn_frame = ttk.Frame(root, padding=(10, 5))
        btn_frame.grid(row=1, column=0, sticky="ew")
        self.create_btn = ttk.Button(btn_frame, text="Create Video", command=self._on_create)
        self.create_btn.pack(side="left")
        self.progress = ttk.Progressbar(btn_frame, mode="indeterminate", length=200)
        self.progress.pack(side="left", padx=10)

        # --- Log output ---
        log_frame = ttk.LabelFrame(root, text="Progress", padding=5)
        log_frame.grid(row=2, column=0, padx=10, pady=(5, 10), sticky="nsew")
        self.log_text = tk.Text(log_frame, height=18, width=90, state="disabled",
                                bg="#1e1e1e", fg="#d4d4d4", font=("Consolas", 9))
        scrollbar = ttk.Scrollbar(log_frame, orient="vertical", command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=scrollbar.set)
        self.log_text.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

    # --- Browse helpers ---
    def _browse_pptx(self):
        path = filedialog.askopenfilename(
            title="Select PowerPoint file",
            filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")],
        )
        if path:
            self.pptx_var.set(path)

    def _browse_outdir(self):
        path = filedialog.askdirectory(title="Select build output directory")
        if path:
            self.outdir_var.set(path)

    def _browse_video(self):
        path = filedialog.asksaveasfilename(
            title="Save video as",
            defaultextension=".mp4",
            filetypes=[("MP4 Video", "*.mp4"), ("All files", "*.*")],
        )
        if path:
            self.video_var.set(path)

    # --- Logging ---
    def _log(self, message):
        """Thread-safe log append."""
        self.root.after(0, self._append_log, message)

    def _append_log(self, message):
        self.log_text.configure(state="normal")
        self.log_text.insert("end", message + "\n")
        self.log_text.see("end")
        self.log_text.configure(state="disabled")

    # --- Run pipeline ---
    def _on_create(self):
        pptx_path = self.pptx_var.get().strip()
        output_dir = self.outdir_var.get().strip()
        output_video = self.video_var.get().strip()
        voice = self.voice_var.get().strip()

        if not pptx_path:
            messagebox.showerror("Error", "Please specify a PowerPoint file.")
            return
        if not os.path.exists(pptx_path):
            if not messagebox.askyesno("Warning", f"File not found:\n{pptx_path}\n\nContinue anyway?"):
                return

        os.makedirs(output_dir, exist_ok=True)

        # Disable button and start progress bar
        self.create_btn.configure(state="disabled")
        self.progress.start(10)

        # Monkey-patch print so pipeline output goes to the log
        import builtins
        original_print = builtins.print

        def log_print(*args, **kwargs):
            message = " ".join(str(a) for a in args)
            self._log(message)
        builtins.print = log_print

        def worker():
            try:
                asyncio.run(run_pipeline(pptx_path, output_dir, output_video, voice, self._log))
                self.root.after(0, lambda: messagebox.showinfo("Complete", f"Video saved to:\n{output_video}"))
            except Exception as e:
                self._log(f"\nERROR: {e}")
                self.root.after(0, lambda: messagebox.showerror("Error", str(e)))
            finally:
                builtins.print = original_print
                self.root.after(0, self._on_done)

        threading.Thread(target=worker, daemon=True).start()

    def _on_done(self):
        self.progress.stop()
        self.create_btn.configure(state="normal")


if __name__ == "__main__":
    root = tk.Tk()
    app = VideoCreatorApp(root)
    root.mainloop()
