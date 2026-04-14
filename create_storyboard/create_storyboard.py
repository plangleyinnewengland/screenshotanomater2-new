from pptx import Presentation
from pptx.util import Inches
import os
import tkinter as tk
from tkinter import filedialog

# Let user pick the screenshots folder
tk.Tk().withdraw()
screenshots_dir = filedialog.askdirectory(title="Select the folder containing screenshots")
if not screenshots_dir:
    print("No folder selected. Exiting.")
    exit()

output_file = os.path.join(screenshots_dir, "storyboard.pptx")

# Get all PNG files, sorted
png_files = sorted([f for f in os.listdir(screenshots_dir) if f.lower().endswith('.png')])

# Create presentation with widescreen dimensions (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Blank slide layout
blank_layout = prs.slide_layouts[6]

for png_file in png_files:
    slide = prs.slides.add_slide(blank_layout)
    
    # Full path to the image
    img_path = os.path.join(screenshots_dir, png_file)
    
    # Add image centered and sized to fit the slide
    # Leave small margins
    left = Inches(0.25)
    top = Inches(0.25)
    width = Inches(12.833)
    height = Inches(7.0)
    
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    
    # Add speaker notes with file path
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = img_path

# Save the presentation
prs.save(output_file)
print(f"Created presentation with {len(png_files)} slides: {output_file}")
